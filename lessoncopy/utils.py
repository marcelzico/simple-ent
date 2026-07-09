from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
import fitz  # PyMuPDF
from lecon.models import Chapter
from .models import Copy, Importer
import os
import json
import re
from io import BytesIO
from zipfile import ZipFile
from django.core.files.base import ContentFile


def iter_block_items(parent):
    """
    Generate a reference to each paragraph and table child within *parent*,
    in document order. Each returned value is an instance of either Table or Paragraph.
    """
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    else:
        raise ValueError("Something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def get_paragraph_list_info(paragraph):
    """Return (is_list_item, level) if paragraph is part of a list."""
    try:
        p = paragraph._element
        numPr = p.xpath('.//w:numPr')
        if numPr:
            ilvl = numPr[0].xpath('.//w:ilvl/@w:val')
            level = int(ilvl[0]) if ilvl else 0
            return True, level
    except Exception as e:
        print(f"Error detecting list info: {e}")
    return False, 0


def extract_docx_to_model(docx_path, chapter_id, importer_instance=None):
    """
    Extract DOCX content and create Copy objects with proper heading levels
    and structured content.
    """
    doc = Document(docx_path)
    chapter = Chapter.objects.get(id=chapter_id)
    
    # First, delete any existing copies for this importer (if re-uploading)
    if importer_instance:
        Copy.objects.filter(chapter=chapter, importer=importer_instance).delete()
    
    # Extract all images first
    image_map = extract_and_map_images(docx_path)
    image_counter = 0
    
    # Track the current heading context
    heading_stack = []  # Stack of (heading_text, heading_level)
    current_content_blocks = []
    last_heading_copy = None
    
    # Process all elements in document order
    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            # Check for images
            images_in_paragraph = extract_images_from_paragraph(block, image_map)
            
            if images_in_paragraph:
                # Save any accumulated content before image
                if current_content_blocks and heading_stack:
                    save_content_blocks(chapter, heading_stack, current_content_blocks, 
                                      importer_instance, last_heading_copy)
                    current_content_blocks = []
                
                # Save each image
                for img_data, img_filename in images_in_paragraph:
                    image_counter += 1
                    django_file = ContentFile(img_data)
                    django_file.name = img_filename
                    
                    caption = extract_image_caption(block) or f"Image {image_counter}"
                    
                    # Create image copy with current heading context
                    if heading_stack:
                        latest_heading = heading_stack[-1]
                        copy = Copy.objects.create(
                            chapter=chapter,
                            heading=latest_heading[0],
                            heading_level=str(latest_heading[1]),
                            image=django_file,
                            caption=caption,
                            importer=importer_instance
                        )
                        last_heading_copy = copy
                    else:
                        copy = Copy.objects.create(
                            chapter=chapter,
                            heading="Introduction",
                            heading_level="1",
                            image=django_file,
                            caption=caption,
                            importer=importer_instance
                        )
                        last_heading_copy = copy
                continue
            
            # Handle headings
            if block.style.name.startswith("Heading"):
                # Save any accumulated content before new heading
                if current_content_blocks and heading_stack:
                    save_content_blocks(chapter, heading_stack, current_content_blocks, 
                                      importer_instance, last_heading_copy)
                    current_content_blocks = []
                
                # Extract heading level
                level_str = block.style.name.replace("Heading", "").strip()
                try:
                    heading_level = int(level_str)
                except ValueError:
                    heading_level = 1
                
                heading_text = block.text.strip()
                
                # Update heading stack
                # Remove any headings with same or higher level
                while heading_stack and heading_stack[-1][1] >= heading_level:
                    heading_stack.pop()
                
                # Add new heading to stack
                heading_stack.append((heading_text, heading_level))
                
                # Create a copy for the heading itself (with no content)
                copy = Copy.objects.create(
                    chapter=chapter,
                    heading=heading_text,
                    heading_level=str(heading_level),
                    content=None,
                    importer=importer_instance
                )
                last_heading_copy = copy
                
            # Handle regular text (paragraphs and list items)
            elif block.text.strip() and not images_in_paragraph:
                text = block.text.strip()
                is_list, list_level = get_paragraph_list_info(block)
                
                content_block = {
                    'type': 'bullet' if is_list else 'paragraph',
                    'level': list_level if is_list else 0,
                    'text': text
                }
                current_content_blocks.append(content_block)
        
        elif isinstance(block, Table):
            # Save any accumulated content before table
            if current_content_blocks and heading_stack:
                save_content_blocks(chapter, heading_stack, current_content_blocks, 
                                  importer_instance, last_heading_copy)
                current_content_blocks = []
            
            # Extract and save table
            table_data = []
            for row in block.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_has_content(table_data) and heading_stack:
                copy = Copy.objects.create(
                    chapter=chapter,
                    heading=heading_stack[-1][0],
                    heading_level=str(heading_stack[-1][1]),
                    content=None,
                    table=table_data,
                    importer=importer_instance
                )
                last_heading_copy = copy
    
    # Save any remaining content
    if current_content_blocks and heading_stack:
        save_content_blocks(chapter, heading_stack, current_content_blocks, 
                          importer_instance, last_heading_copy)


def save_content_blocks(chapter, heading_stack, content_blocks, importer_instance, last_copy=None):
    """
    Save accumulated content blocks under the current heading.
    Each heading gets its own Copy with the content blocks stored in JSON.
    """
    if not content_blocks or not heading_stack:
        return
    
    # Get the current heading (last in stack)
    current_heading = heading_stack[-1][0]
    current_level = heading_stack[-1][1]
    
    # Clean and optimize blocks
    cleaned_blocks = []
    for block in content_blocks:
        if block.get('text') and block['text'].strip():
            cleaned_block = block.copy()
            cleaned_block['text'] = block['text'].strip()
            cleaned_blocks.append(cleaned_block)
    
    if not cleaned_blocks:
        return
    
    content_data = {
        'type': 'section',
        'heading': current_heading,
        'level': current_level,
        'children': cleaned_blocks
    }
    
    # Check if we should update the last copy or create a new one
    if last_copy and last_copy.heading == current_heading and last_copy.content is None:
        # Update the heading copy with content
        last_copy.content = content_data
        last_copy.save()
    else:
        # Create new copy
        Copy.objects.create(
            chapter=chapter,
            heading=current_heading,
            heading_level=str(current_level),
            content=content_data,
            importer=importer_instance
        )


def extract_pptx_to_model(pptx_path, chapter_id, importer_instance=None):
    """
    Extract PowerPoint content and create copies with proper structure.
    PowerPoint slides become separate copies with their content as structured JSON.
    """
    chapter = Chapter.objects.get(id=chapter_id)
    prs = Presentation(pptx_path)
    
    # Delete existing copies for this importer
    if importer_instance:
        Copy.objects.filter(chapter=chapter, importer=importer_instance).delete()
    
    for i, slide in enumerate(prs.slides):
        slide_title = None
        content_blocks = []
        
        # First pass: try to find title
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and not slide_title:
                if shape == slide.shapes.title:
                    slide_title = shape.text.strip()
                    break
        
        # Second pass: collect all text content
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # Skip if this is the title (already captured)
                if slide_title and text == slide_title:
                    continue
                
                # Check if it's a bullet point
                if paragraph.level > 0 or text.startswith(('•', '-', '●', '·')):
                    # Remove bullet character
                    clean_text = re.sub(r'^[•\-●·\s]+', '', text).strip()
                    content_blocks.append({
                        'type': 'bullet',
                        'level': paragraph.level,
                        'text': clean_text
                    })
                else:
                    content_blocks.append({
                        'type': 'paragraph',
                        'text': text
                    })
        
        if content_blocks:
            content_data = {
                'type': 'section',
                'heading': slide_title or f"Slide {i+1}",
                'level': 2,
                'children': content_blocks
            }
        else:
            content_data = None
        
        Copy.objects.create(
            chapter=chapter,
            heading=slide_title or f"Slide {i+1}",
            heading_level="2",
            content=content_data,
            importer=importer_instance
        )


def extract_pdf_to_model(pdf_path, chapter_id, importer_instance=None):
    """
    Extract PDF content and create copies.
    Each page becomes a separate copy with its content as structured JSON.
    """
    chapter = Chapter.objects.get(id=chapter_id)
    doc = fitz.open(pdf_path)
    
    # Delete existing copies for this importer
    if importer_instance:
        Copy.objects.filter(chapter=chapter, importer=importer_instance).delete()
    
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            # Try to detect title (first line or all caps line)
            lines = text.split('\n')
            title = None
            content_blocks = []
            
            for j, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                # First line might be title
                if j == 0 and len(line) < 100:
                    title = line
                # Line in all caps might be heading
                elif line.isupper() and len(line) < 100:
                    # Create a new heading copy
                    if content_blocks:
                        # Save previous content
                        content_data = {
                            'type': 'section',
                            'heading': title or f"Page {i+1}",
                            'level': 2,
                            'children': content_blocks
                        }
                        Copy.objects.create(
                            chapter=chapter,
                            heading=title or f"Page {i+1}",
                            heading_level="2",
                            content=content_data,
                            importer=importer_instance
                        )
                        content_blocks = []
                     
                    title = line
                    Copy.objects.create(
                        chapter=chapter,
                        heading=line,
                        heading_level="2",
                        content=None,
                        importer=importer_instance
                    )
                # Check for bullet points
                elif line.startswith(('•', '-', '●', '·', '*')) or re.match(r'^\d+\.', line):
                    clean_text = re.sub(r'^[•\-●·*\d\.\s]+', '', line).strip()
                    content_blocks.append({
                        'type': 'bullet',
                        'level': 0,
                        'text': clean_text
                    })
                else:
                    content_blocks.append({
                        'type': 'paragraph',
                        'text': line
                    })
            
            # Save remaining content
            if content_blocks:
                content_data = {
                    'type': 'section',
                    'heading': title or f"Page {i+1}",
                    'level': 2,
                    'children': content_blocks
                }
                Copy.objects.create(
                    chapter=chapter,
                    heading=title or f"Page {i+1}",
                    heading_level="2",
                    content=content_data,
                    importer=importer_instance
                )


def extract_and_map_images(docx_path):
    """Extract all images from docx and create a mapping"""
    image_map = {}
    try:
        with ZipFile(docx_path, 'r') as docx_zip:
            # Parse relationships
            relationships = {}
            if 'word/_rels/document.xml.rels' in docx_zip.namelist():
                rels_content = docx_zip.read('word/_rels/document.xml.rels').decode('utf-8')
                pattern = r'Relationship Id="([^"]+)"[^>]+Target="media/([^"]+)"'
                matches = re.findall(pattern, rels_content)
                for r_id, filename in matches:
                    relationships[r_id] = filename
            
            # Extract images
            image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
            for image_file in image_files:
                filename = os.path.basename(image_file)
                image_data = docx_zip.read(image_file)
                
                # Find relationship ID
                found = False
                for r_id, rel_filename in relationships.items():
                    if rel_filename == filename:
                        image_map[r_id] = {
                            'data': image_data,
                            'filename': filename
                        }
                        found = True
                        break
                
                if not found:
                    image_map[filename] = {
                        'data': image_data,
                        'filename': filename
                    }
    except Exception as e:
        print(f"Error extracting images: {e}")
    
    return image_map


def extract_images_from_paragraph(paragraph, image_map):
    """Extract images from a paragraph"""
    images = []
    try:
        for run in paragraph.runs:
            graphic = run._element.xpath('.//a:graphic')
            if graphic:
                blip = graphic[0].xpath('.//a:blip')
                if blip:
                    r_id = blip[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if r_id and r_id in image_map:
                        img_data = image_map[r_id]['data']
                        img_filename = image_map[r_id]['filename']
                        images.append((img_data, img_filename))
    except Exception as e:
        print(f"Error extracting images: {e}")
    
    return images


def extract_image_caption(paragraph):
    """Extract image caption from paragraph"""
    text = paragraph.text.strip()
    
    caption_patterns = [
        r'Figure\s+\d+[:\-]\s*(.+)',
        r'Fig\.\s*\d+[:\-]\s*(.+)',
        r'Image\s+\d+[:\-]\s*(.+)',
        r'Table\s+\d+[:\-]\s*(.+)',
    ]
    
    for pattern in caption_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return match.group(1).strip()
    
    if text and len(text) < 200:
        return text
    
    return None


def table_has_content(table_data):
    """Check if table has any non-empty content"""
    for row in table_data:
        for cell in row:
            if cell and cell.strip():
                return True
    return False


def prepare_table_data(data, padding=2, max_col_width=30, header=False):
    """Process data into a format suitable for template rendering."""
    if not data:
        return None

    cols = max(len(row) if isinstance(row, (list, tuple)) else 1 for row in data)
    processed_data = []
    col_widths = [0] * cols
    
    for row in data:
        if not isinstance(row, (list, tuple)):
            row = [row]
        
        processed_row = []
        for i, cell in enumerate(row[:cols]):
            cell_str = str(cell) if cell is not None else ""
            lines = cell_str.split('\n')
            processed_row.append(lines)
            col_widths[i] = max(col_widths[i], *[min(len(line), max_col_width) for line in lines])
        
        processed_row += [[] for _ in range(cols - len(row))]
        processed_data.append(processed_row)
    
    col_widths = [w + padding * 2 for w in col_widths]
    
    return {
        'rows': processed_data,
        'col_widths': col_widths,
        'header': header,
        'max_col_width': max_col_width
    }

